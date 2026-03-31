"""
Document Processor — handles extraction from any file type.
Supports: PDF, DOCX, TXT, CSV, XLSX, PPTX, MD, JSON, HTML
"""
import os
import io
import json
import chardet
import pandas as pd
from pathlib import Path
from typing import List, Tuple
from loguru import logger

# LangChain text splitter
from langchain_text_splitters import RecursiveCharacterTextSplitter

# File-specific loaders
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from config import get_settings

settings = get_settings()


class DocumentProcessor:
    def __init__(self):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )

    def extract_text(self, file_path: str, filename: str) -> Tuple[str, dict]:
        """Extract raw text and metadata from any supported file."""
        ext = Path(filename).suffix.lower()
        logger.info(f"Extracting text from: {filename} (type: {ext})")

        metadata = {
            "filename": filename,
            "file_type": ext,
            "file_size_kb": round(os.path.getsize(file_path) / 1024, 2),
        }

        try:
            if ext == ".pdf":
                text, pages = self._extract_pdf(file_path)
                metadata["pages"] = pages
            elif ext in (".docx", ".doc"):
                text = self._extract_docx(file_path)
            elif ext in (".txt", ".md", ".rst", ".log"):
                text = self._extract_text_file(file_path)
            elif ext == ".csv":
                text = self._extract_csv(file_path)
            elif ext in (".xlsx", ".xls"):
                text = self._extract_excel(file_path)
            elif ext == ".pptx":
                text, slides = self._extract_pptx(file_path)
                metadata["slides"] = slides
            elif ext == ".json":
                text = self._extract_json(file_path)
            elif ext in (".html", ".htm"):
                text = self._extract_html(file_path)
            else:
                # Fallback: try reading as plain text
                text = self._extract_text_file(file_path)

            metadata["char_count"] = len(text)
            metadata["word_count"] = len(text.split())
            logger.success(f"Extracted {len(text)} chars from {filename}")
            return text, metadata

        except Exception as e:
            logger.error(f"Failed to extract {filename}: {e}")
            raise ValueError(f"Could not process file '{filename}': {str(e)}")

    def chunk_text(self, text: str, metadata: dict) -> List[dict]:
        """Split text into overlapping chunks with metadata."""
        chunks = self.splitter.split_text(text)
        logger.info(f"Split into {len(chunks)} chunks")

        enriched_chunks = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():  # skip empty chunks
                enriched_chunks.append({
                    "text": chunk.strip(),
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    **metadata,
                })
        return enriched_chunks

    # ── Private extractors ─────────────────────────────────────────────────

    def _extract_pdf(self, path: str) -> Tuple[str, int]:
        reader = PdfReader(path)
        pages = len(reader.pages)
        text_parts = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"[Page {i+1}]\n{page_text}")
        return "\n\n".join(text_parts), pages

    def _extract_docx(self, path: str) -> str:
        doc = DocxDocument(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading structure
                if para.style.name.startswith("Heading"):
                    parts.append(f"\n## {para.text.strip()}\n")
                else:
                    parts.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_text_file(self, path: str) -> str:
        with open(path, "rb") as f:
            raw = f.read()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding") or "utf-8"
        return raw.decode(encoding, errors="replace")

    def _extract_csv(self, path: str) -> str:
        try:
            df = pd.read_csv(path, encoding="utf-8")
        except UnicodeDecodeError:
            df = pd.read_csv(path, encoding="latin-1")
        # Convert to readable text: column descriptions + rows
        text_parts = [f"Columns: {', '.join(df.columns.tolist())}"]
        text_parts.append(f"Total rows: {len(df)}")
        text_parts.append("")
        # Include all rows as text
        for _, row in df.iterrows():
            row_text = " | ".join(f"{col}: {val}" for col, val in row.items() if pd.notna(val))
            text_parts.append(row_text)
        return "\n".join(text_parts)

    def _extract_excel(self, path: str) -> str:
        xl = pd.ExcelFile(path)
        text_parts = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text_parts.append(f"=== Sheet: {sheet_name} ===")
            text_parts.append(f"Columns: {', '.join(df.columns.astype(str).tolist())}")
            for _, row in df.iterrows():
                row_text = " | ".join(
                    f"{col}: {val}" for col, val in row.items() if pd.notna(val)
                )
                text_parts.append(row_text)
            text_parts.append("")
        return "\n".join(text_parts)

    def _extract_pptx(self, path: str) -> Tuple[str, int]:
        prs = Presentation(path)
        slides = len(prs.slides)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts), slides

    def _extract_json(self, path: str) -> str:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)

    def _extract_html(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        # Simple HTML tag stripping
        import re
        clean = re.sub(r"<[^>]+>", " ", content)
        clean = re.sub(r"\s+", " ", clean).strip()
        return clean
